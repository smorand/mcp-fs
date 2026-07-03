"""Lightweight document text extraction to Markdown (pure Python, lazy imports).

Ideas adapted from ``../docling-scripts`` but WITHOUT its docling / torch / VLM
stack: we call the underlying format libraries directly (pypdf, python-docx,
python-pptx, openpyxl, html2text, Pillow) so mcp-fs stays deployable, no native
runtime. Audio and video are out of scope.

Extraction is CPU bound and synchronous; callers must run it off the event loop
(``asyncio.to_thread``), like every other blocking call in this codebase.
"""

from __future__ import annotations

import csv
import io
from dataclasses import dataclass, field
from typing import Any

_TEXT_EXTS = frozenset({".txt", ".md", ".markdown", ".rst", ".log", ".text"})
_FENCED_EXTS = frozenset({".json", ".yaml", ".yml", ".xml", ".toml", ".ini", ".env"})
_IMAGE_EXTS = frozenset({".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"})
_AV_EXTS = frozenset({".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".mp4", ".mkv", ".mov", ".avi", ".webm", ".wmv"})
_TABLE_ROW_CAP = 400


@dataclass(slots=True)
class ExtractResult:
    """Outcome of a document extraction."""

    fmt: str
    text: str = ""
    truncated: bool = False
    meta: dict[str, Any] = field(default_factory=dict)
    note: str = ""

    def as_dict(self, path: str) -> dict[str, Any]:
        return {
            "path": path,
            "format": self.fmt,
            "text": self.text,
            "chars": len(self.text),
            "truncated": self.truncated,
            "meta": self.meta,
            "note": self.note,
        }


class UnsupportedDocument(Exception):
    """Raised for formats we explicitly do not extract (audio/video)."""


def _ext(filename: str) -> str:
    return ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""


def extract(data: bytes, filename: str, *, max_chars: int = 200_000, ocr: bool = True) -> ExtractResult:
    """Extract Markdown/plain text from ``data`` chosen by ``filename`` extension."""
    ext = _ext(filename)
    if ext in _AV_EXTS:
        raise UnsupportedDocument(f"audio/video is out of scope for extraction: {ext}")
    if ext == ".pdf":
        result = _pdf(data)
    elif ext == ".docx":
        result = _docx(data)
    elif ext in {".pptx", ".pptm", ".potx", ".ppsx"}:
        result = _pptx(data)
    elif ext in {".xlsx", ".xlsm"}:
        result = _xlsx(data)
    elif ext in {".html", ".htm"}:
        result = _html(data)
    elif ext == ".csv":
        result = _csv(data)
    elif ext in _IMAGE_EXTS:
        result = _image(data, ocr=ocr)
    elif ext in _FENCED_EXTS:
        result = _fenced(data, ext)
    elif ext in _TEXT_EXTS or ext == "":
        result = ExtractResult("text", _decode(data))
    else:
        # Unknown extension: best-effort decode as UTF-8 text.
        result = ExtractResult("text", _decode(data), note=f"unknown extension {ext}; decoded as text")
    if len(result.text) > max_chars:
        result.text = result.text[:max_chars]
        result.truncated = True
    return result


def _decode(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _md_table(rows: list[list[str]]) -> str:
    """Render rows (first row = header) as a GitHub Markdown table."""
    rows = [[(c if c is not None else "").replace("\n", " ").replace("|", "\\|").strip() for c in row] for row in rows]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header, *body = rows
    out = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    out += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(out)


def _pdf(data: bytes) -> ExtractResult:
    from pypdf import PdfReader  # noqa: PLC0415 - lazy, heavy import

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            parts.append(f"\n\n---\n*[Page {index}]*\n\n{text.strip()}")
    body = "".join(parts).strip()
    meta: dict[str, Any] = {"pages": len(reader.pages)}
    info = reader.metadata
    if info and info.title:
        meta["title"] = str(info.title)
    note = "" if body else "no extractable text layer (scanned PDF?); OCR is not enabled in this build"
    return ExtractResult("pdf", body, meta=meta, note=note)


def _docx(data: bytes) -> ExtractResult:
    import docx  # noqa: PLC0415

    document = docx.Document(io.BytesIO(data))
    lines: list[str] = []
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name if para.style else "") or ""
        if style.lower().startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit())
            hashes = "#" * min(int(level) if level else 1, 6)
            lines.append(f"{hashes} {text}")
        elif style.lower().startswith("list"):
            lines.append(f"- {text}")
        else:
            lines.append(text)
    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows][:_TABLE_ROW_CAP]
        if rows:
            lines.append("")
            lines.append(_md_table(rows))
    return ExtractResult("docx", "\n\n".join(lines).strip(), meta={"paragraphs": len(document.paragraphs)})


def _pptx(data: bytes) -> ExtractResult:
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {index}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = ("".join(run.text for run in para.runs) or para.text).strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows][:_TABLE_ROW_CAP]
                if rows:
                    lines.append(_md_table(rows))
        notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        if notes:
            lines.append(f"> Notes: {notes}")
    return ExtractResult("pptx", "\n\n".join(lines).strip(), meta={"slides": len(prs.slides)})


def _xlsx(data: bytes) -> ExtractResult:
    from openpyxl import load_workbook  # noqa: PLC0415

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        rows = [["" if value is None else str(value) for value in row] for row in sheet.iter_rows(values_only=True)]
        rows = [r for r in rows if any(c.strip() for c in r)][:_TABLE_ROW_CAP]
        if not rows:
            continue
        lines.append(f"## Sheet: {sheet.title}")
        lines.append(_md_table(rows))
    workbook.close()
    return ExtractResult("xlsx", "\n\n".join(lines).strip(), meta={"sheets": len(workbook.sheetnames)})


def _html(data: bytes) -> ExtractResult:
    import html2text  # noqa: PLC0415

    converter = html2text.HTML2Text()
    converter.body_width = 0
    converter.ignore_images = True
    text = converter.handle(_decode(data)).strip()
    return ExtractResult("html", text)


def _csv(data: bytes) -> ExtractResult:
    reader = csv.reader(io.StringIO(_decode(data)))
    rows = list(reader)[:_TABLE_ROW_CAP]
    return ExtractResult("csv", _md_table(rows), meta={"rows": len(rows)})


def _fenced(data: bytes, ext: str) -> ExtractResult:
    lang = ext.lstrip(".")
    return ExtractResult(lang, f"```{lang}\n{_decode(data)}\n```")


def _image(data: bytes, *, ocr: bool) -> ExtractResult:
    meta: dict[str, Any] = {}
    try:
        from PIL import Image  # noqa: PLC0415

        with Image.open(io.BytesIO(data)) as image:
            meta = {"width": image.width, "height": image.height, "mode": image.mode}
    except Exception:
        meta = {}
    if ocr:
        text = _try_ocr(data)
        if text:
            return ExtractResult("image", text, meta=meta, note="text recovered via Tesseract OCR (CPU, degraded mode)")
    note = (
        "image: no OCR text. Visual understanding needs a multimodal model (not wired into mcp-fs); "
        "install Tesseract for CPU OCR."
    )
    return ExtractResult("image", "", meta=meta, note=note)


def _try_ocr(data: bytes) -> str:
    """Best-effort CPU OCR. Silently returns '' if Tesseract/pytesseract is absent."""
    try:
        import pytesseract  # noqa: PLC0415
        from PIL import Image  # noqa: PLC0415

        with Image.open(io.BytesIO(data)) as image:
            return str(pytesseract.image_to_string(image)).strip()
    except Exception:
        return ""
