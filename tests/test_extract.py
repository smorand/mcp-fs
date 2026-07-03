"""Document extraction + Markdown->docx rendering (pure, no live stack)."""

from __future__ import annotations

import io

import pytest

from mcp_fs.docx_writer import markdown_to_docx
from mcp_fs.extract import UnsupportedDocument, extract


def test_docx_roundtrip_headings_table_and_bold() -> None:
    md = (
        "# Titre\n\nUn **paragraphe** clef.\n\n- point A\n- point B\n\n| Nom | Rôle |\n| --- | --- |\n| Seb | Archi |\n"
    )
    data = markdown_to_docx(md, title="Synthèse")
    result = extract(data, "out.docx")
    assert result.fmt == "docx"
    assert "Titre" in result.text
    assert "point A" in result.text
    assert "Seb" in result.text and "Archi" in result.text  # table survived


def test_pptx_slides_and_text() -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide de test"
    buffer = io.BytesIO()
    prs.save(buffer)
    result = extract(buffer.getvalue(), "deck.pptx")
    assert result.fmt == "pptx"
    assert result.meta["slides"] == 1
    assert "Slide de test" in result.text
    assert "## Slide 1" in result.text


def test_xlsx_sheet_to_markdown_table() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Ville", "Pop"])
    ws.append(["Kenitra", 431282])
    buffer = io.BytesIO()
    wb.save(buffer)
    result = extract(buffer.getvalue(), "cities.xlsx")
    assert result.fmt == "xlsx"
    assert "## Sheet: Data" in result.text
    assert "Kenitra" in result.text and "431282" in result.text


def test_csv_and_html_and_text() -> None:
    assert "| a | b |" in extract(b"a,b\n1,2\n", "x.csv").text
    html = extract(b"<h1>Hi</h1><p>Body</p>", "x.html")
    assert html.fmt == "html" and "# Hi" in html.text and "Body" in html.text
    assert extract(b"plain notes", "notes.md").text == "plain notes"


def test_json_is_fenced() -> None:
    result = extract(b'{"a": 1}', "config.json")
    assert result.text.startswith("```json") and '"a": 1' in result.text


def test_truncation_flag() -> None:
    result = extract(b"x" * 100, "big.txt", max_chars=10)
    assert result.truncated and len(result.text) == 10


def test_image_without_ocr_returns_note_and_meta() -> None:
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (12, 8), "white").save(buffer, format="PNG")
    result = extract(buffer.getvalue(), "pic.png", ocr=False)
    assert result.fmt == "image"
    assert result.meta["width"] == 12 and result.meta["height"] == 8
    assert result.text == "" and "multimodal" in result.note


def test_audio_video_rejected() -> None:
    for name in ("song.mp3", "clip.mp4", "call.m4a"):
        with pytest.raises(UnsupportedDocument):
            extract(b"", name)


def test_docx_writer_numbered_list_and_italic() -> None:
    data = markdown_to_docx("1. premier\n2. second\n\n*emphase* simple\n")
    result = extract(data, "x.docx")
    assert "premier" in result.text and "second" in result.text and "emphase" in result.text
